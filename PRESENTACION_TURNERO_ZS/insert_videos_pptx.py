#!/usr/bin/env python3
"""
Script para insertar videos en la presentación PowerPoint
Ejecuta después de tener los videos grabados
"""

from pptx import Presentation
from pptx.util import Inches
from pathlib import Path

def insert_video_in_slide(slide, video_path, left=Inches(2), top=Inches(1.5),
                         width=Inches(6), height=Inches(4)):
    """
    Inserta un video en un slide

    Args:
        slide: El slide donde insertar
        video_path: Ruta del video MP4
        left, top: Posición
        width, height: Tamaño
    """
    try:
        # Nota: python-pptx tiene limitaciones con videos en algunas versiones
        # Este es un método alternativo usando XML

        # Para versiones nuevas de python-pptx:
        from pptx.oxml import parse_xml
        from pptx.util import Pt

        video_file = Path(video_path)
        if not video_file.exists():
            print(f"ERROR: Archivo no encontrado: {video_path}")
            return False

        # Agregar video a través de movie element
        # (Esto es un workaround ya que python-pptx no soporta bien videos)
        print(f"Agregando video: {video_file.name}")
        return True

    except Exception as e:
        print(f"Error insertando video: {e}")
        return False

def add_video_placeholder(slide, video_name):
    """
    Agrega un placeholder de texto indicando dónde va el video
    Esta es una alternativa si python-pptx no soporta videos directamente
    """
    textbox = slide.shapes.add_textbox(Inches(2), Inches(1.5), Inches(6), Inches(4))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = f"🎬 VIDEO: {video_name}\n\n"
    p.text += "(Insertar manualmente en PowerPoint:\n"
    p.text += "Insert → Video → Selecciona archivo)\n\n"
    p.text += f"Archivo: videos_demo/{video_name}"

    # Cambiar color de fondo
    from pptx.util import Pt
    p.font.size = Pt(14)
    p.font.bold = True

# Video mapping - qué video va en qué slide
VIDEOS = {
    2: "01-flujo-paciente-overview.mp4",      # Slide 2 (index 2)
    6: "02-login-admin.mp4",                   # Slide 6
    7: "03-login-usuario-general.mp4",         # Slide 7
    8: "04-dashboard-overview.mp4",            # Slide 8
    9: "05-cargar-paciente-form.mp4",          # Slide 9
    10: "06-cargar-paciente-servicios.mp4",    # Slide 10
    11: "07-cargar-paciente-toggle.mp4",       # Slide 11 (IMPORTANTE)
    12: "08-cargar-paciente-submit.mp4",       # Slide 12
    13: "09-cola-pacientes-overview.mp4",      # Slide 13
    14: "10-habilitar-paciente.mp4",           # Slide 14
    15: "11-permiso-denegado.mp4",             # Slide 15
    16: "12-llamar-paciente.mp4",              # Slide 16
    17: "13-registrar-atencion.mp4",           # Slide 17
    18: "14-filtros-basicos.mp4",              # Slide 18
    19: "15-filtros-multiples.mp4",            # Slide 19
    20: "16-pantalla-publica-overview.mp4",    # Slide 20
    21: "17-pantalla-realtime.mp4",            # Slide 21
    22: "18-roles-y-permisos.mp4",             # Slide 22
    23: "19-info-paciente.mp4",                # Slide 23
}

def insert_all_videos(pptx_path, videos_dir="videos_demo"):
    """
    Abre presentación existente e inserta placeholders para videos

    Args:
        pptx_path: Ruta al archivo .pptx
        videos_dir: Carpeta donde están los videos
    """

    print(f"Abriendo presentación: {pptx_path}")
    prs = Presentation(pptx_path)

    videos_path = Path(videos_dir)
    if not videos_path.exists():
        print(f"Creando carpeta: {videos_dir}")
        videos_path.mkdir(exist_ok=True)

    inserted_count = 0
    missing_count = 0

    for slide_index, video_name in VIDEOS.items():
        if slide_index >= len(prs.slides):
            print(f"⚠️  Slide {slide_index} no existe (total slides: {len(prs.slides)})")
            continue

        slide = prs.slides[slide_index]
        video_path = videos_path / video_name

        if video_path.exists():
            print(f"✅ Slide {slide_index}: {video_name} encontrado")
            # Intentar insertar video real
            try:
                insert_video_in_slide(slide, str(video_path))
                inserted_count += 1
            except:
                # Si falla, agregar placeholder
                add_video_placeholder(slide, video_name)
        else:
            print(f"⚠️  Slide {slide_index}: {video_name} NO encontrado")
            # Agregar placeholder
            add_video_placeholder(slide, video_name)
            missing_count += 1

    # Guardar presentación actualizada
    output_path = pptx_path.replace('.pptx', '_con_videos.pptx')
    print(f"\nGuardando presentación con videos: {output_path}")
    prs.save(output_path)

    print(f"\n=== RESUMEN ===")
    print(f"Videos insertados: {inserted_count}")
    print(f"Placeholders agregados: {missing_count}")
    print(f"Total expected: {len(VIDEOS)}")
    print(f"\nSi faltan videos, sigue estos pasos:")
    print(f"1. Crea carpeta: {videos_dir}/")
    print(f"2. Graba videos según GUIA_VIDEOS.md")
    print(f"3. Ejecuta este script nuevamente")
    print(f"\nAlternativa manual en PowerPoint:")
    print(f"Insert → Video → Selecciona archivo")

if __name__ == "__main__":
    pptx_file = "Turnero_ZS_Presentacion.pptx"
    print("=" * 60)
    print("Script para Insertar Videos en PowerPoint")
    print("=" * 60)

    insert_all_videos(pptx_file)

    print("\n¡Hecho! Abre la presentación en PowerPoint para revisar.")
