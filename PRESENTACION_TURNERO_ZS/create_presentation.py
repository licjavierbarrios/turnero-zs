#!/usr/bin/env python3
"""
Generador de Presentación PowerPoint para Turnero ZS
Crea automáticamente una presentación con los slides de la demo
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Crear presentación
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Colores corporativos
COLOR_PRIMARY = RGBColor(59, 130, 246)      # Azul
COLOR_SECONDARY = RGBColor(16, 185, 129)   # Verde
COLOR_ACCENT = RGBColor(249, 115, 22)      # Naranja
COLOR_DARK = RGBColor(15, 23, 42)          # Gris oscuro
COLOR_LIGHT = RGBColor(241, 245, 249)      # Gris claro

def add_title_slide(prs, title, subtitle):
    """Agregar slide de título"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(60)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    subtitle_frame.word_wrap = True

def add_content_slide(prs, title, content_points, video_info=None):
    """Agregar slide de contenido con puntos"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Fondo blanco
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Línea de color en top
    line = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.1))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_PRIMARY
    line.line.color.rgb = COLOR_PRIMARY

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLOR_DARK

    # Contenido
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, point in enumerate(content_points):
        if i > 0:
            text_frame.add_paragraph()

        p = text_frame.paragraphs[i]
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_DARK
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        p.level = 0

    # Video info si existe
    if video_info:
        video_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.2), Inches(8.6), Inches(0.9))
        video_frame = video_box.text_frame
        video_frame.text = f"📹 Video: {video_info}"
        video_frame.paragraphs[0].font.size = Pt(14)
        video_frame.paragraphs[0].font.italic = True
        video_frame.paragraphs[0].font.color.rgb = COLOR_SECONDARY

def add_two_column_slide(prs, title, left_content, right_content):
    """Agregar slide con dos columnas"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Fondo blanco
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Línea de color en top
    line = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.1))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_PRIMARY
    line.line.color.rgb = COLOR_PRIMARY

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLOR_DARK

    # Columna izquierda
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(5.8))
    left_text = left_box.text_frame
    left_text.word_wrap = True

    for i, point in enumerate(left_content):
        if i > 0:
            left_text.add_paragraph()
        p = left_text.paragraphs[i]
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_DARK
        p.space_before = Pt(4)
        p.space_after = Pt(4)

    # Columna derecha
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(5.8))
    right_text = right_box.text_frame
    right_text.word_wrap = True

    for i, point in enumerate(right_content):
        if i > 0:
            right_text.add_paragraph()
        p = right_text.paragraphs[i]
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_DARK
        p.space_before = Pt(4)
        p.space_after = Pt(4)

# ============================================================================
# SLIDES
# ============================================================================

# Slide 1: Portada
add_title_slide(prs, "Turnero ZS",
    "Sistema de Gestión de Turnos\nPara Centros de Salud Argentinos")

# Slide 2: Problema y Contexto
add_content_slide(prs, "Problema y Contexto",
    [
        "❌ Demoras prolongadas en atención a pacientes",
        "❌ Desorganización en las colas de espera",
        "❌ Información fragmentada entre sistemas",
        "❌ Falta de visibilidad en tiempo real",
        "",
        "✅ Solución: Sistema integrado de gestión de turnos",
        "✅ Contexto: Sistema de salud argentino (CAPS/Hospitales)",
        "✅ Objetivo: Reducir tiempos de espera y mejorar experiencia"
    ])

# Slide 3: Características Principales
add_content_slide(prs, "Características Principales",
    [
        "✅ Gestión de turnos en tiempo real",
        "✅ Cola de pacientes diaria con múltiples estados",
        "✅ Pantalla pública de avance (Realtime)",
        "✅ Control de profesionales y consultorios",
        "✅ Sistema de roles y permisos granulares",
        "✅ Múltiples instituciones en una plataforma",
        "✅ Toggle para cargar pacientes habilitados o pendientes",
        "✅ Audio TTS para llamada de pacientes"
    ])

# Slide 4: Flujo del Paciente
add_content_slide(prs, "Flujo del Paciente en el Sistema",
    [
        "📋 PENDIENTE → ✅ DISPONIBLE → 📢 LLAMADO → ✓ ATENDIDO",
        "",
        "• Pendiente: Paciente registrado, requiere habilitación",
        "  (Solo el creador puede habilitar)",
        "",
        "• Disponible: Paciente habilitado, listo para ser llamado",
        "",
        "• Llamado: Paciente siendo llamado (con audio TTS)",
        "",
        "• Atendido: Consulta completada (fin del proceso)"
    ],
    "01-flujo-paciente-overview.mp4")

# Slide 5: Login y Autenticación
add_content_slide(prs, "Demo 1: Login y Autenticación",
    [
        "🔐 Autenticación segura con Supabase Auth",
        "",
        "• Ingresar credenciales (email + password)",
        "• Sistema reconoce múltiples instituciones asignadas",
        "• Seleccionar institución de trabajo",
        "• Acceder al dashboard según rol",
        "",
        "Videos:",
        "  • 02-login-admin.mp4",
        "  • 03-login-usuario-general.mp4"
    ])

# Slide 6: Dashboard Principal
add_content_slide(prs, "Demo 2: Dashboard Principal",
    [
        "📊 Interfaz principal del sistema",
        "",
        "Información visible:",
        "• Cantidad de pacientes por estado (Pendiente, Disponible, etc.)",
        "• Filtros avanzados (Servicio, Profesional, Consultorio, Estado)",
        "• Lista de pacientes en tiempo real",
        "• Botón para cargar nuevo paciente",
        "• Información de cada paciente (hora de carga, profesional, etc.)",
    ],
    "04-dashboard-overview.mp4")

# Slide 7: Cargar Paciente - Parte 1
add_content_slide(prs, "Demo 3: Cargar Nuevo Paciente (Parte 1)",
    [
        "📝 Abrir diálogo 'Cargar Nuevo Paciente'",
        "",
        "Pasos:",
        "1. Clic en botón '+ Cargar Paciente'",
        "2. Modal scroll aparece (se puede desplazar si hay mucho contenido)",
        "3. Ingresar Nombre Completo",
        "4. Ingresar DNI",
        "5. Validación automática de campos",
    ],
    "05-cargar-paciente-form.mp4")

# Slide 8: Cargar Paciente - Parte 2
add_content_slide(prs, "Demo 4: Cargar Nuevo Paciente (Parte 2)",
    [
        "🏥 Seleccionar servicios/profesionales",
        "",
        "Pasos:",
        "1. Checkboxes múltiples para servicios",
        "2. Checkboxes múltiples para profesionales asignados hoy",
        "3. Contador de seleccionados",
        "4. Puede seleccionar múltiples opciones",
        "5. Modal scroll permite ver muchas opciones",
    ],
    "06-cargar-paciente-servicios.mp4")

# Slide 9: Toggle de Estado Inicial (FEATURE NUEVA)
add_content_slide(prs, "Demo 5: Toggle Estado Inicial (⭐ NUEVA FEATURE)",
    [
        "🎚️ Control de estado inicial del paciente",
        "",
        "Opciones:",
        "⟳ Pendiente (defecto, ámbar)",
        "   → Requiere habilitación posterior",
        "   → Solo el creador puede habilitar",
        "",
        "✓ Disponible (verde)",
        "   → Inmediatamente disponible para atención",
        "   → Ya está habilitado",
        "",
        "💡 Permite flexibilidad en el flujo de carga",
    ],
    "07-cargar-paciente-toggle.mp4")

# Slide 10: Confirmar Carga
add_content_slide(prs, "Demo 6: Confirmar Carga de Paciente",
    [
        "✅ Finalizar carga del paciente",
        "",
        "Pasos:",
        "1. Clic en botón 'Cargar Paciente'",
        "2. Modal se cierra automáticamente",
        "3. Paciente aparece en la cola",
        "4. Estado según selección (Pendiente o Disponible)",
        "5. Hora de carga se registra automáticamente",
        "6. Optimistic UI: aparece inmediatamente",
    ],
    "08-cargar-paciente-submit.mp4")

# Slide 11: Cola de Pacientes
add_content_slide(prs, "Demo 7: Gestión de Pacientes en Cola",
    [
        "📋 Vista completa de la cola del día",
        "",
        "Información visible por paciente:",
        "• Número de orden (001, 002, 003, etc.)",
        "• Nombre y DNI del paciente",
        "• Servicio solicitado",
        "• Estado con código de color",
        "• Profesional y consultorio asignado",
        "• Hora de carga (con 🕐 icon)",
        "• Botones de acción según estado",
    ],
    "09-cola-pacientes-overview.mp4")

# Slide 12: Habilitar Paciente (Permisos)
add_content_slide(prs, "Demo 8: Habilitar Paciente (Control de Permisos)",
    [
        "🔒 Solo el admin que cargó el paciente puede habilitarlo",
        "",
        "Casos:",
        "✅ Creador del paciente:",
        "   → Ve botón 'Habilitar' activo",
        "   → Puede cambiar de Pendiente a Disponible",
        "",
        "❌ Otro usuario:",
        "   → Ve botón 'Habilitar' deshabilitado",
        "   → Muestra icono 🔒 (candado)",
        "   → Explicación en tooltip",
        "",
        "💡 Seguridad: Solo quien carga controla habilitación"
    ],
    "10-habilitar-paciente.mp4")

# Slide 13: Permiso Denegado
add_content_slide(prs, "Demo 9: Control de Permisos en Acción",
    [
        "🚫 Visualizar restricción de permisos",
        "",
        "Escenario:",
        "• Paciente cargado por Admin A",
        "• Admin B intenta habilitar",
        "• Sistema muestra: Botón deshabilitado con 🔒",
        "",
        "Beneficios:",
        "✅ Responsabilidad clara (quién cargó, quién habilita)",
        "✅ Previene cambios no autorizados",
        "✅ Trazabilidad del proceso",
        "✅ Seguridad del flujo",
    ],
    "11-permiso-denegado.mp4")

# Slide 14: Llamar Paciente
add_content_slide(prs, "Demo 10: Llamar Paciente (Audio TTS)",
    [
        "📢 Sistema de llamada con audio en español",
        "",
        "Proceso:",
        "1. Seleccionar paciente en estado 'Disponible'",
        "2. Clic en botón 'Llamar'",
        "3. Audio TTS anuncia: 'Paciente [nombre], consultorio [número]'",
        "4. Duración: ~11 segundos (dos anuncios)",
        "5. Estado cambia a 'Llamado'",
        "",
        "💡 Diferenciador: Audio generado en tiempo real",
        "💡 Accesibilidad: Ayuda a pacientes con discapacidad visual",
    ],
    "12-llamar-paciente.mp4")

# Slide 15: Registrar Atención
add_content_slide(prs, "Demo 11: Registrar Atención Completada",
    [
        "✓ Marcar paciente como atendido",
        "",
        "Proceso:",
        "1. Paciente en estado 'Llamado' (en consulta)",
        "2. Clic en botón 'Registrar Atención'",
        "3. Estado cambia a 'Atendido'",
        "4. Timestamp automático de fin",
        "5. Paciente completa su flujo",
        "",
        "Datos registrados:",
        "• Hora de carga",
        "• Hora de habilitación (si aplica)",
        "• Hora de llamada",
        "• Hora de atención completada"
    ],
    "13-registrar-atencion.mp4")

# Slide 16: Filtros - Básicos
add_content_slide(prs, "Demo 12: Filtros Avanzados (Parte 1)",
    [
        "🔍 Filtrado por criterios individuales",
        "",
        "Opciones de filtro:",
        "• Por Servicio (Cardiología, Pediatría, etc.)",
        "• Por Profesional (Nombre del doctor)",
        "• Por Consultorio (A, B, C, etc.)",
        "• Por Estado (Pendiente, Disponible, Llamado, Atendido)",
        "",
        "Interacción:",
        "1. Seleccionar filtro en dropdown",
        "2. Cola se actualiza inmediatamente",
        "3. Mostrar cantidad de resultados",
        "4. Botón 'Limpiar filtros' para resetear"
    ],
    "14-filtros-basicos.mp4")

# Slide 17: Filtros - Múltiples
add_content_slide(prs, "Demo 13: Filtros Avanzados (Parte 2)",
    [
        "🔍 Combinación de múltiples filtros",
        "",
        "Ejemplos:",
        "• Filtrar: Servicio=Cardiología + Estado=Disponible",
        "• Filtrar: Profesional=Dr. García + Estado=Pendiente",
        "• Filtrar: Consultorio=A + Servicio=Pediatría",
        "",
        "Resultados:",
        "• Actualización en tiempo real",
        "• Contador de pacientes que cumplen criterios",
        "• Todos los filtros se aplican simultáneamente",
        "• Limpiar todo con un clic",
    ],
    "15-filtros-multiples.mp4")

# Slide 18: Pantalla Pública - Overview
add_content_slide(prs, "Demo 14: Pantalla Pública (Parte 1)",
    [
        "📺 Visualización pública de la cola para pacientes",
        "",
        "Características:",
        "• URL diferente: /pantalla/[institution-id]",
        "• No requiere login (o solo rol 'pantalla')",
        "• Información clara y legible",
        "• Diseño atractivo y simple",
        "• Actualización automática en tiempo real",
        "• Responsive (funciona en TV, tablet, mobile)",
        "",
        "Información visible:",
        "• Próximo paciente a ser atendido",
        "• Servicio y profesional",
        "• Consultorio asignado"
    ],
    "16-pantalla-publica-overview.mp4")

# Slide 19: Pantalla Pública - Realtime
add_content_slide(prs, "Demo 15: Sincronización en Tiempo Real",
    [
        "⚡ Actualización instantánea sin recargar",
        "",
        "Flujo demostrativo:",
        "1. Pantalla pública abierta en una TV/monitor",
        "2. Admin carga paciente en dashboard",
        "3. Paciente aparece INMEDIATAMENTE en pantalla (Supabase Realtime)",
        "4. Admin habilita paciente",
        "5. Estado se actualiza en pantalla",
        "6. Admin llamar paciente",
        "7. Cambio visible en tiempo real",
        "",
        "💡 Tecnología: Supabase Realtime Channels",
        "💡 Diferenciador: No requiere polling o refresco"
    ],
    "17-pantalla-realtime.mp4")

# Slide 20: Roles y Permisos
add_content_slide(prs, "Demo 16: Sistema de Roles y Permisos",
    [
        "👥 Diferentes vistas según rol del usuario",
        "",
        "Roles implementados:",
        "👤 Admin: Acceso completo (todos los servicios)",
        "👤 Administrativo: Cargar pacientes, habilitar",
        "👤 Médico: Solo sus servicios y pacientes",
        "👤 Enfermería: Auxiliar del administrativo",
        "👤 Pantalla: Solo lectura de cola pública",
        "",
        "Demostración:",
        "• Login con diferentes usuarios",
        "• Mostrar interfaz diferente por rol",
        "• Explicar permisos de cada rol"
    ],
    "18-roles-y-permisos.mp4")

# Slide 21: Información de Paciente
add_content_slide(prs, "Demo 17: Detalles de Paciente",
    [
        "ℹ️ Información completa de cada paciente",
        "",
        "Datos visibles:",
        "• Nombre completo del paciente",
        "• DNI",
        "• Número de orden (001, 002, etc.)",
        "• Servicio seleccionado",
        "• Profesional asignado",
        "• Consultorio asignado",
        "• Hora de carga (con 🕐 icon)",
        "• Estado actual",
        "• Timestamps de transiciones (si aplica)",
    ],
    "19-info-paciente.mp4")

# Slide 22: Stack Tecnológico
add_two_column_slide(prs, "Stack Tecnológico",
    [
        "Frontend:",
        "• Next.js 15.5.2",
        "• React 19",
        "• TypeScript",
        "• Tailwind CSS 4",
        "• shadcn/ui 3",
        "",
        "Testing:",
        "• Vitest",
        "• React Testing Library",
        "• 152 tests passing"
    ],
    [
        "Backend:",
        "• Supabase",
        "• PostgreSQL",
        "• Supabase Auth",
        "• Supabase Realtime",
        "",
        "DevOps:",
        "• GitHub Actions",
        "• Vercel Deployment",
        "• Row Level Security (RLS)",
        "• Multi-tenancy"
    ])

# Slide 23: Métricas y Resultados
add_content_slide(prs, "Métricas y Resultados Alcanzados",
    [
        "✅ Objetivos del MVP:",
        "  • Reducir tiempo de espera en 25-40%",
        "  • Disminuir absentismo en 10-20%",
        "  • Alcanzar ≥85% ocupación de horarios",
        "  • Mantener ≥95% trazabilidad completa",
        "",
        "✅ Características implementadas:",
        "  • Gestión de cola diaria en tiempo real",
        "  • Sistema de llamada con audio TTS",
        "  • Pantalla pública con Realtime",
        "  • Control granular de permisos",
        "  • 152 tests automatizados",
        "  • Sin errores de typecheck/lint"
    ])

# Slide 24: Ventajas Competitivas
add_content_slide(prs, "Ventajas Competitivas",
    [
        "🚀 Sistema integrado (sin cambios en HSI)",
        "⚡ Tiempo real (sin recargas)",
        "📱 Responsive (desktop, tablet, mobile)",
        "🔒 Seguro (RLS, autenticación, permisos)",
        "🌐 Multi-tenancy (múltiples instituciones)",
        "♿ Accesible (WCAG compliant)",
        "📊 Escalable (PostgreSQL + Supabase)",
        "🎯 Intuitivo (UI clara y lógica)",
        "🎚️ Toggle para estado inicial (flexibilidad)",
        "📢 Audio TTS en español (diferenciador)"
    ])

# Slide 25: Roadmap Futuro
add_two_column_slide(prs, "Roadmap Futuro",
    [
        "Corto Plazo (1-2 meses):",
        "✓ Integración HSI",
        "✓ Reportes avanzados",
        "✓ Notificaciones push",
        "",
        "Mediano Plazo (3-6 meses):",
        "✓ App móvil para pacientes",
        "✓ Confirmación por SMS",
        "✓ Asignación automática"
    ],
    [
        "Largo Plazo (6+ meses):",
        "✓ Predicción de demora (ML)",
        "✓ Gestor de camas",
        "✓ Sistema de emergencia",
        "",
        "Opcionales:",
        "✓ Integración con PACS",
        "✓ Teleconsulta",
        "✓ Analítica avanzada"
    ])

# Slide 26: Conclusión y CTA
add_content_slide(prs, "Conclusión y Próximos Pasos",
    [
        "✅ Sistema completo de gestión de turnos",
        "✅ Mejora significativa en experiencia del paciente",
        "✅ Fácil integración con institutos existentes",
        "",
        "📈 Impacto esperado:",
        "• Reducción de tiempos de espera",
        "• Mayor satisfacción de pacientes",
        "• Optimización de recursos",
        "• Mejor trazabilidad de procesos",
        "",
        "🎯 Próximos pasos:",
        "• Feedback de instituciones piloto",
        "• Refinamiento según necesidades",
        "• Rollout a más instituciones"
    ])

# Slide 27: Contacto y Links
add_title_slide(prs, "¡Gracias!",
    "GitHub: github.com/licjavierbarrios/turnero-zs\n\n" +
    "Email: licjavierbarrios@gmail.com\n\n" +
    "Demo: [URL en producción]")

# Guardar presentación
output_path = "Turnero_ZS_Presentacion.pptx"
prs.save(output_path)
print("[OK] Presentacion creada exitosamente: " + output_path)
print("[INFO] Total de slides: " + str(len(prs.slides)))
print("[INFO] Videos necesarios: 19")
print("[INFO] Duracion total de videos: 7-8 minutos")
